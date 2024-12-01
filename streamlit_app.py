import streamlit as st
import google.generativeai as genai
import os
from io import StringIO
import PyPDF2

# Configure the API key securely from Streamlit's secrets
# Ensure that GOOGLE_API_KEY is added in secrets.toml (for local) or Streamlit Cloud Secrets
genai.configure(api_key=st.secrets["GOOGLE_API_KEY"])

# Streamlit App UI
st.title("Pre-Sales Content Optimization Tool")
st.write("""
    This tool uses AI to optimize your pre-sales content like pitch decks, demo scripts, and proposals.
    It suggests improvements for keyword optimization, content effectiveness, and visual enhancements.
""")

# Upload content or enter text manually
content_input = st.text_area("Enter your content (pitch deck, demo script, proposal, etc.):", "")

# File upload for content
uploaded_file = st.file_uploader("Or upload your content (PDF, Word, PPT)", type=["pdf", "docx", "pptx"])

# If a file is uploaded, process the content
if uploaded_file is not None:
    if uploaded_file.type == "application/pdf":
        # Convert PDF to text (using a simple method, or you could use a library like PyPDF2)
        import PyPDF2
        pdf_reader = PyPDF2.PdfReader(uploaded_file)
        text = ""
        for page in pdf_reader.pages:
            text += page.extract_text()
        content_input = text
    elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        import docx
        doc = docx.Document(uploaded_file)
        content_input = "\n".join([para.text for para in doc.paragraphs])
    elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        from pptx import Presentation
        prs = Presentation(uploaded_file)
        content_input = "\n".join([slide.shapes.title.text for slide in prs.slides if slide.shapes.title])

# Button to analyze content and provide suggestions
if st.button("Analyze Content"):
    if content_input:
        try:
            # Load and configure the model
            model = genai.GenerativeModel('gemini-1.5-flash')
            
            # Generate AI suggestions for content improvement
            prompt = f"Analyze this pre-sales content for keyword optimization, content effectiveness, and visual enhancement suggestions: {content_input}"
            response = model.generate_content(prompt)
            
            # Display response in Streamlit
            st.write("Suggestions for Optimizing Your Content:")
            st.write(response.text)
            
            # Advanced features for visual and content suggestions
            st.write("""
                ### Visual Enhancements:
                Consider these tips for improving your pitch deck visuals:
                - Use high-quality, relevant images for better engagement.
                - Consistent color schemes and fonts help make your deck look professional.
                - Avoid overloading slides with too much text; prioritize key points and visuals.
                - Include charts and infographics to break down complex information.
            """)
        except Exception as e:
            st.error(f"Error generating suggestions: {e}")
    else:
        st.warning("Please enter content or upload a file to analyze.")
        
# Industry-specific keyword optimization
st.sidebar.title("Industry-Specific Keyword Optimization")
industry_selection = st.sidebar.selectbox("Select Industry", ["Tech", "Finance", "Healthcare", "Retail"])

# Ensure model is loaded for keyword optimization
model = genai.GenerativeModel('gemini-1.5-flash')  # Ensure model is loaded before use

if industry_selection:
    st.sidebar.write(f"Optimizing for industry: {industry_selection}")
    # Generate AI-driven keyword optimization based on the selected industry
    industry_keywords_prompt = f"Suggest the most relevant keywords for the {industry_selection} industry to optimize pre-sales content."
    try:
        # Generate keywords based on selected industry
        response_keywords = model.generate_content(industry_keywords_prompt)
        st.sidebar.write("Suggested Keywords:")
        st.sidebar.write(response_keywords.text)
    except Exception as e:
        st.sidebar.error(f"Error generating industry-specific keywords: {e}")

# Content Version Comparison
st.sidebar.title("Compare Content Versions")
st.sidebar.write("""
    Upload a previous version of your content (e.g., previous pitch deck) to compare against the current version.
    This helps track improvements in content effectiveness over time.
""")

# File upload for old content version
uploaded_old_file = st.sidebar.file_uploader("Upload previous version of content (PDF, Word, PPT)", type=["pdf", "docx", "pptx"])

if uploaded_old_file is not None:
    if uploaded_old_file.type == "application/pdf":
        # Convert PDF to text
        pdf_reader = PyPDF2.PdfReader(uploaded_old_file)
        old_text = ""
        for page in pdf_reader.pages:
            old_text += page.extract_text()
        old_content_input = old_text
    elif uploaded_old_file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        doc = docx.Document(uploaded_old_file)
        old_content_input = "\n".join([para.text for para in doc.paragraphs])
    elif uploaded_old_file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        prs = Presentation(uploaded_old_file)
        old_content_input = "\n".join([slide.shapes.title.text for slide in prs.slides if slide.shapes.title])

    # Button to analyze old content version
    if st.sidebar.button("Analyze Previous Content Version"):
        if old_content_input:
            try:
                # Generate AI suggestions for old content
                old_prompt = f"Analyze this previous version of pre-sales content and suggest improvements: {old_content_input}"
                old_response = model.generate_content(old_prompt)
                
                st.sidebar.write("Suggestions for Improving Previous Content Version:")
                st.sidebar.write(old_response.text)
            except Exception as e:
                st.sidebar.error(f"Error generating suggestions for old content: {e}")

# Additional analytics for content
st.sidebar.title("Content Engagement Analytics")
st.sidebar.write("""
    Track engagement for specific content sections, slides, or proposals. This feature would require integration with analytics tools, but here are some generic ideas:
    - Track the most-viewed slides or sections in your pitch decks.
    - Get feedback on sections that received less attention.
    - Analyze the performance of your proposals based on past results.
""")
